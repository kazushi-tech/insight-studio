"""PPTX export service — generates editable PowerPoint from ReviewResult.

This is a **prototype behind feature flag** (CR-A4.3).
Enable via FEATURE_PPTX_EXPORT=1 environment variable.
Does NOT block HTML/PDF export pipeline.

Requires: python-pptx (optional dependency).
"""

from __future__ import annotations

import logging
import os
from typing import Optional

from ...schemas.review_result import ReviewResult

logger = logging.getLogger("market-lens")

FEATURE_FLAG = "FEATURE_PPTX_EXPORT"


class PptxExportError(Exception):
    """Raised when PPTX generation fails."""


class PptxFeatureDisabledError(PptxExportError):
    """Raised when the PPTX feature flag is not enabled."""


def is_pptx_enabled() -> bool:
    """Check whether PPTX export feature flag is on."""
    return os.environ.get(FEATURE_FLAG, "").strip() in ("1", "true", "yes")


def export_pptx(
    result: ReviewResult,
    *,
    title: str = "クリエイティブレビュー サマリー",
    subtitle: str = "",
    review_date: Optional[str] = None,
) -> bytes:
    """Generate an editable PPTX from a ReviewResult.

    Args:
        result: The structured review output.
        title: Presentation title.
        subtitle: Subtitle (brand/campaign).
        review_date: Date string (YYYY-MM-DD). Defaults to today.

    Returns:
        PPTX file content as bytes.

    Raises:
        PptxFeatureDisabledError: Feature flag not enabled.
        PptxExportError: python-pptx unavailable or generation fails.
    """
    if not is_pptx_enabled():
        raise PptxFeatureDisabledError(
            f"PPTX export is disabled. Set {FEATURE_FLAG}=1 to enable."
        )

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches, Pt  # type: ignore[import-untyped]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    except ImportError:
        raise PptxExportError(
            "python-pptx is not installed. Install it with: pip install python-pptx"
        )

    from datetime import date
    import io

    if review_date is None:
        review_date = date.today().isoformat()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_textbox(
        slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
        title, Pt(32), bold=True, alignment=PP_ALIGN.CENTER,
    )
    sub_text = subtitle or ""
    if review_date:
        sub_text = f"{sub_text}  |  {review_date}" if sub_text else review_date
    _add_textbox(
        slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
        sub_text, Pt(18), alignment=PP_ALIGN.CENTER,
    )

    # --- Slide 2: Summary + Good Points ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
        "サマリー", Pt(24), bold=True,
    )
    _add_textbox(
        slide, Inches(0.5), Inches(1.0), Inches(12), Inches(1.5),
        result.summary, Pt(14),
    )
    _add_textbox(
        slide, Inches(0.5), Inches(2.8), Inches(12), Inches(0.6),
        "良い点 (Good Points)", Pt(20), bold=True,
    )
    good_text = "\n".join(
        f"  {p.point} — {p.reason}" for p in result.good_points
    )
    _add_textbox(
        slide, Inches(0.5), Inches(3.5), Inches(12), Inches(3.5),
        good_text, Pt(12),
    )

    # --- Slide 3: Keep + Improvements ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(5.5), Inches(0.6),
        "守るべき点 (Keep)", Pt(20), bold=True,
    )
    keep_text = "\n".join(
        f"  {p.point} — {p.reason}" for p in result.keep_as_is
    )
    _add_textbox(
        slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(5.5),
        keep_text, Pt(12),
    )
    _add_textbox(
        slide, Inches(6.8), Inches(0.3), Inches(6), Inches(0.6),
        "改善提案 (Improvements)", Pt(20), bold=True,
    )
    improve_text = "\n".join(
        f"  {p.point}\n    理由: {p.reason}\n    Action: {p.action}"
        for p in result.improvements
    )
    _add_textbox(
        slide, Inches(6.8), Inches(1.0), Inches(6), Inches(5.5),
        improve_text, Pt(12),
    )

    # --- Slide 4: Test Ideas + Evidence ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(5.5), Inches(0.6),
        "テスト案 (Test Ideas)", Pt(20), bold=True,
    )
    test_text = "\n".join(
        f"  {t.hypothesis}\n    変数: {t.variable} / 期待効果: {t.expected_impact}"
        for t in result.test_ideas
    )
    _add_textbox(
        slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(5.5),
        test_text, Pt(12),
    )
    _add_textbox(
        slide, Inches(6.8), Inches(0.3), Inches(6), Inches(0.6),
        "根拠・出典 (Evidence)", Pt(20), bold=True,
    )
    evidence_text = "\n".join(
        f"  [{e.evidence_type.value}] {e.evidence_source}\n    {e.evidence_text}"
        for e in result.evidence
    )
    _add_textbox(
        slide, Inches(6.8), Inches(1.0), Inches(6), Inches(5.5),
        evidence_text, Pt(12),
    )

    # Serialize to bytes
    buf = io.BytesIO()
    prs.save(buf)
    pdf_bytes = buf.getvalue()
    logger.info("PPTX export: generated %d bytes (%d slides)", len(pdf_bytes), len(prs.slides))
    return pdf_bytes


def _add_textbox(slide, left, top, width, height, text, font_size,
                 *, bold=False, alignment=None):
    """Helper to add an editable textbox to a slide."""
    from pptx.util import Pt  # type: ignore[import-untyped]

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    if alignment is not None:
        p.alignment = alignment
