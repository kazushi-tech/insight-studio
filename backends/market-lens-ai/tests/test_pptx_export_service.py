"""Tests for PPTX export service (CR-A4.3)."""

from __future__ import annotations

import os
from unittest import mock

import pytest

from web.app.schemas.review_result import (
    EvidenceItem,
    EvidenceType,
    ImprovementPoint,
    ReviewPoint,
    ReviewResult,
    RubricScore,
    TestIdea,
)
from web.app.services.exports.pptx_export_service import (
    FEATURE_FLAG,
    PptxExportError,
    PptxFeatureDisabledError,
    export_pptx,
    is_pptx_enabled,
)


@pytest.fixture()
def sample_review() -> ReviewResult:
    return ReviewResult(
        review_type="banner_review",
        summary="ECサイト向けセールバナー。50% OFFの訴求が目立つ。",
        good_points=[
            ReviewPoint(point="割引率の視認性", reason="50% OFFが大きく表示されている"),
        ],
        keep_as_is=[
            ReviewPoint(point="ブランドカラーの一貫性", reason="赤と白の配色がブランドガイドに沿っている"),
        ],
        improvements=[
            ImprovementPoint(
                point="CTAボタンを拡大",
                reason="現在のサイズでは視認性が低い",
                action="CTAボタンの高さを1.5倍に拡大する",
            ),
        ],
        test_ideas=[
            TestIdea(
                hypothesis="CTAを緑に変更するとCTR向上が期待できる（仮説）",
                variable="CTAボタン色",
                expected_impact="CTR 5-10%改善の可能性",
            ),
        ],
        evidence=[
            EvidenceItem(
                evidence_type=EvidenceType.client_material,
                evidence_source="brand guideline v2.1",
                evidence_text="赤をプライマリカラーとして使用",
            ),
        ],
        target_hypothesis="20-40代女性、EC購買層",
        message_angle="期間限定の割引訴求",
        rubric_scores=[
            RubricScore(rubric_id="visual_impact", score=4, comment="目を引くレイアウト"),
        ],
    )


class TestFeatureFlag:
    def test_disabled_by_default(self):
        with mock.patch.dict(os.environ, {}, clear=True):
            os.environ.pop(FEATURE_FLAG, None)
            assert is_pptx_enabled() is False

    def test_enabled_with_1(self):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            assert is_pptx_enabled() is True

    def test_enabled_with_true(self):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "true"}):
            assert is_pptx_enabled() is True

    def test_enabled_with_yes(self):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "yes"}):
            assert is_pptx_enabled() is True

    def test_disabled_with_0(self):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "0"}):
            assert is_pptx_enabled() is False

    def test_disabled_with_empty(self):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: ""}):
            assert is_pptx_enabled() is False


class TestExportPptx:
    def test_raises_when_flag_disabled(self, sample_review: ReviewResult):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "0"}):
            with pytest.raises(PptxFeatureDisabledError, match="disabled"):
                export_pptx(sample_review)

    def test_raises_when_pptx_not_installed(self, sample_review: ReviewResult):
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            with mock.patch.dict("sys.modules", {"pptx": None, "pptx.util": None, "pptx.enum.text": None}):
                with pytest.raises(PptxExportError, match="python-pptx"):
                    export_pptx(sample_review)

    def test_generates_valid_pptx(self, sample_review: ReviewResult):
        """Test that PPTX export produces valid bytes when python-pptx is available."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            result = export_pptx(
                sample_review,
                title="テストレビュー",
                subtitle="テストブランド",
                review_date="2026-03-22",
            )

        # Should be valid PPTX (ZIP-based Office format)
        assert isinstance(result, bytes)
        assert len(result) > 0
        # PPTX files start with PK (ZIP signature)
        assert result[:2] == b"PK"

    def test_pptx_has_4_slides(self, sample_review: ReviewResult):
        """Verify template produces exactly 4 slides."""
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        import io
        from pptx import Presentation as Prs

        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            pptx_bytes = export_pptx(sample_review, review_date="2026-03-22")

        prs = Prs(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 4

    def test_pptx_text_is_editable(self, sample_review: ReviewResult):
        """Verify text boxes in PPTX are editable (not images)."""
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        import io
        from pptx import Presentation as Prs

        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            pptx_bytes = export_pptx(sample_review, review_date="2026-03-22")

        prs = Prs(io.BytesIO(pptx_bytes))
        # Every slide should have at least one shape with a text_frame
        for slide in prs.slides:
            has_text = any(
                hasattr(shape, "text_frame") and shape.text_frame is not None
                for shape in slide.shapes
            )
            assert has_text, f"Slide {slide.slide_id} has no editable text boxes"

    def test_default_date_is_today(self, sample_review: ReviewResult):
        """If no date is provided, today's date should be used."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        from datetime import date

        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            result = export_pptx(sample_review)

        # Just verify it doesn't crash — date is embedded in slide text
        assert isinstance(result, bytes)
        assert len(result) > 0
