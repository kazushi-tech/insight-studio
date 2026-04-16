"""PPTX edit roundtrip test — proves generated PPTX is editable (M5.4 / WP-B4)."""

from __future__ import annotations

import io
import os
from unittest import mock

import pytest

from web.app.schemas.review_result import (
    EvidenceItem,
    EvidenceType,
    ImprovementPoint,
    ReviewPoint,
    ReviewResult,
    RubricScore,
    TestIdea,
)
from web.app.services.exports.pptx_export_service import (
    FEATURE_FLAG,
    export_pptx,
)


@pytest.fixture()
def sample_review() -> ReviewResult:
    return ReviewResult(
        review_type="banner_review",
        summary="ラウンドトリップテスト用レビュー",
        good_points=[
            ReviewPoint(point="良い点", reason="理由"),
        ],
        keep_as_is=[
            ReviewPoint(point="維持点", reason="理由"),
        ],
        improvements=[
            ImprovementPoint(point="改善点", reason="理由", action="アクション"),
        ],
        test_ideas=[
            TestIdea(hypothesis="仮説", variable="変数", expected_impact="効果"),
        ],
        evidence=[
            EvidenceItem(
                evidence_type=EvidenceType.client_material,
                evidence_source="出典",
                evidence_text="根拠テキスト",
            ),
        ],
        target_hypothesis="ターゲット",
        message_angle="メッセージ",
        rubric_scores=[
            RubricScore(rubric_id="test", score=4, comment="コメント"),
        ],
    )


class TestPptxEditRoundtrip:
    """Generate PPTX, edit it, re-save, and verify the edit persisted."""

    def test_edit_title_roundtrip(self, sample_review, tmp_path):
        """Modify title text in slide 1, re-save, verify modification persists."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Step 1: Generate PPTX
        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            pptx_bytes = export_pptx(
                sample_review,
                title="元のタイトル",
                subtitle="テストブランド",
                review_date="2026-03-23",
            )

        # Step 2: Open the generated PPTX
        prs = Presentation(io.BytesIO(pptx_bytes))
        slide1 = prs.slides[0]

        # Find the title text box (first shape with the original title)
        title_shape = None
        for shape in slide1.shapes:
            if hasattr(shape, "text_frame") and "元のタイトル" in shape.text_frame.text:
                title_shape = shape
                break
        assert title_shape is not None, "Could not find title shape in slide 1"

        # Step 3: Modify the title text
        new_title = "変更後のタイトル"
        title_shape.text_frame.paragraphs[0].text = new_title

        # Step 4: Re-save to a new buffer
        save_path = tmp_path / "edited.pptx"
        prs.save(str(save_path))

        # Step 5: Re-open and verify the modification persisted
        prs2 = Presentation(str(save_path))
        slide1_reopened = prs2.slides[0]
        found_new_title = False
        for shape in slide1_reopened.shapes:
            if hasattr(shape, "text_frame") and new_title in shape.text_frame.text:
                found_new_title = True
                break
        assert found_new_title, f"Modified title '{new_title}' not found after re-open"

    def test_edit_summary_roundtrip(self, sample_review, tmp_path):
        """Modify summary text in slide 2, re-save, verify it persists."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        with mock.patch.dict(os.environ, {FEATURE_FLAG: "1"}):
            pptx_bytes = export_pptx(sample_review, review_date="2026-03-23")

        prs = Presentation(io.BytesIO(pptx_bytes))
        slide2 = prs.slides[1]

        # Find the summary shape
        summary_shape = None
        for shape in slide2.shapes:
            if hasattr(shape, "text_frame") and "ラウンドトリップテスト" in shape.text_frame.text:
                summary_shape = shape
                break
        assert summary_shape is not None, "Could not find summary shape in slide 2"

        # Modify
        new_summary = "編集されたサマリーテキスト"
        summary_shape.text_frame.paragraphs[0].text = new_summary

        # Re-save
        save_path = tmp_path / "edited_summary.pptx"
        prs.save(str(save_path))

        # Re-open and verify
        prs2 = Presentation(str(save_path))
        found = any(
            hasattr(s, "text_frame") and new_summary in s.text_frame.text
            for s in prs2.slides[1].shapes
        )
        assert found, f"Modified summary '{new_summary}' not found after re-open"
